from pptx import Presentation
from pptx.util import Inches
import os

# Criação da apresentação
ppt = Presentation()

# Slide 1 - Título (layout mais organizado e simples)
slide_1 = ppt.slides.add_slide(ppt.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]

title.text = "Projeto de Inteligência Artificial - Jogo da Velha"
subtitle.text = "Lucas Josué Schneider - Engenharia de Software"

# Slide 2 - Introdução
slide_2 = ppt.slides.add_slide(ppt.slide_layouts[1])
title_2 = slide_2.shapes.title
body_2 = slide_2.shapes.placeholders[1]

title_2.text = "Introdução"
body_2.text = ("Este projeto tem como objetivo desenvolver um jogo da velha com inteligência artificial "
               "aplicando técnicas de aprendizado de máquina. O foco está na implementação de árvores de decisão "
               "ou algoritmos de regressão para prever os movimentos dos jogadores e estratégias de jogo.")

# Slide 3 - O Jogo da Velha (explicando o jogo e incluindo imagem)
slide_3 = ppt.slides.add_slide(ppt.slide_layouts[5])  # Layout com imagem
title_3 = slide_3.shapes.title
title_3.text = "O Jogo da Velha"

# Explicação sobre o jogo
text_box = slide_3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(3))
text_frame = text_box.text_frame
text_frame.text = ("O Jogo da Velha é um jogo de dois jogadores que se alternam em marcar os espaços de um tabuleiro "
                   "3x3 com 'X' ou 'O'. O objetivo é formar uma linha com três símbolos iguais, seja na horizontal, "
                   "vertical ou diagonal.")

# Adicionado imagem do tabuleiro de Jogo da Velha
img_path = "data/tic_tac_toe_image.png"  # Caminho para a imagem do Jogo da Velha
slide_3.shapes.add_picture(img_path, Inches(5), Inches(1.5), width=Inches(3.5))

# Slide 4 - Algoritmos de Aprendizado de Máquina
slide_4 = ppt.slides.add_slide(ppt.slide_layouts[1])
title_4 = slide_4.shapes.title
body_4 = slide_4.shapes.placeholders[1]

title_4.text = "Algoritmos de Aprendizado de Máquina"
body_4.text = ("O modelo de IA foi treinado usando os seguintes algoritmos:\n"
               "- Árvore de Decisão\n"
               "- Regressão\n"
               "Os dados do jogo foram utilizados para treinar o modelo a prever os movimentos com base "
               "no histórico das partidas.")

# Slide 5 - Coleta de Dados
slide_5 = ppt.slides.add_slide(ppt.slide_layouts[1])
title_5 = slide_5.shapes.title
body_5 = slide_5.shapes.placeholders[1]

title_5.text = "Coleta de Dados"
body_5.text = ("O estado do jogo e as ações dos jogadores foram registrados durante 100 partidas de IA vs IA. "
               "Esses dados foram usados para treinar o modelo de IA e ajustar seus hiperparâmetros para obter o melhor desempenho.")

# Slide 6 - Resultados
slide_6 = ppt.slides.add_slide(ppt.slide_layouts[1])
title_6 = slide_6.shapes.title
body_6 = slide_6.shapes.placeholders[1]

title_6.text = "Resultados"
body_6.text = ("Os resultados das 100 partidas foram analisados para avaliar a acurácia do modelo. Gráficos de distribuição "
               "de vitórias, empates e derrotas foram gerados, mostrando a eficácia da IA nas decisões de jogo.")

# Slide 7 - Conclusão
slide_7 = ppt.slides.add_slide(ppt.slide_layouts[1])
title_7 = slide_7.shapes.title
body_7 = slide_7.shapes.placeholders[1]

title_7.text = "Conclusão"
body_7.text = ("O uso de algoritmos de aprendizado de máquina, como árvores de decisão e regressão, mostrou-se eficaz para prever "
               "movimentos em um jogo simples como o jogo da velha. A IA foi capaz de melhorar suas estratégias ao longo das partidas, "
               "e os ajustes nos hiperparâmetros foram fundamentais para o desempenho do modelo.")

# Salvando o arquivo
ppt_file_path = "data/Jogo_da_Velha_IA_Lucas.pptx"
ppt.save(ppt_file_path)

ppt_file_path
